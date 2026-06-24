import re
from collections import OrderedDict
from pathlib import Path
from shutil import which
from subprocess import run
from tempfile import NamedTemporaryFile, TemporaryDirectory

import boto3
from langchain_community.document_loaders import Docx2txtLoader, PyPDFLoader, TextLoader
from langchain_core.documents import Document
from pptx import Presentation


def load_document_from_s3(bucket: str, key: str, *, region_name: str) -> list[Document]:
    """Parse one S3 object into LangChain Document objects.

    The object is downloaded through boto3, then parsed by a file-type-specific
    LangChain loader. This avoids the heavy OCR/layout dependency stack used by
    unstructured for normal text PDFs.
    """
    suffix = Path(key).suffix.lower()
    s3_client = boto3.client("s3", region_name=region_name)

    with NamedTemporaryFile(suffix=suffix, delete=False) as temp_file:
        temp_path = Path(temp_file.name)
        s3_client.download_fileobj(bucket, key, temp_file)

    try:
        loader = _loader_for_path(temp_path, suffix)
        documents = loader.load()
    finally:
        temp_path.unlink(missing_ok=True)

    for document in documents:
        document.page_content = _clean_text(document.page_content)
        document.metadata.update(
            {
                "s3_bucket": bucket,
                "s3_key": key,
                "source": f"s3://{bucket}/{key}",
            }
        )

    return documents


def _loader_for_path(path: Path, suffix: str):
    if suffix == ".pdf":
        return PyPDFLoader(str(path))
    if suffix == ".docx":
        return Docx2txtLoader(str(path))
    if suffix == ".doc":
        return BinaryStringDocumentLoader(path, source_format="doc")
    if suffix == ".pptx":
        return PowerPointLoader(path)
    if suffix == ".ppt":
        return LegacyPowerPointLoader(path)
    if suffix in {".txt", ".md", ".csv"}:
        return TextLoader(str(path), encoding="utf-8")

    raise ValueError(f"Unsupported document type '{suffix}' for {path.name}.")


class PowerPointLoader:
    """Extract text from a PPTX file, one LangChain document per slide."""

    def __init__(self, path: Path) -> None:
        self.path = path

    def load(self) -> list[Document]:
        presentation = Presentation(str(self.path))
        documents: list[Document] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            text_parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)

                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text for cell in row.cells if cell.text]
                        if cells:
                            text_parts.append(" | ".join(cells))

            content = "\n".join(part.strip() for part in text_parts if part.strip())
            if content:
                documents.append(
                    Document(
                        page_content=content,
                        metadata={"slide": slide_index, "total_slides": len(presentation.slides)},
                    )
                )

        return documents


class LegacyPowerPointLoader:
    """Extract text from legacy PPT files.

    LibreOffice gives better results when available. On small EC2 workers where
    it is not installed, fall back to extracting readable strings from the binary
    PowerPoint file.
    """

    def __init__(self, path: Path) -> None:
        self.path = path

    def load(self) -> list[Document]:
        if not which("soffice"):
            return self._load_from_binary_strings()

        with TemporaryDirectory() as temp_dir:
            result = run(
                [
                    "soffice",
                    "--headless",
                    "--convert-to",
                    "pptx",
                    "--outdir",
                    temp_dir,
                    str(self.path),
                ],
                capture_output=True,
                text=True,
                check=False,
            )
            converted_path = Path(temp_dir) / f"{self.path.stem}.pptx"
            if result.returncode != 0 or not converted_path.exists():
                raise RuntimeError(
                    "Unable to parse legacy .ppt file. Install LibreOffice on the worker "
                    "or upload .pptx files. "
                    f"LibreOffice output: {result.stderr or result.stdout}"
                )

            return PowerPointLoader(converted_path).load()

    def _load_from_binary_strings(self) -> list[Document]:
        return BinaryStringDocumentLoader(self.path, source_format="ppt").load()


class BinaryStringDocumentLoader:
    """Best-effort text extraction for legacy binary Office files."""

    def __init__(self, path: Path, *, source_format: str) -> None:
        self.path = path
        self.source_format = source_format

    def load(self) -> list[Document]:
        data = self.path.read_bytes()
        candidates: list[str] = []

        for match in re.finditer(rb"(?:[\x20-\x7e]\x00){4,}", data):
            candidates.append(match.group(0).decode("utf-16le", errors="ignore"))

        for match in re.finditer(rb"[\x20-\x7e]{8,}", data):
            candidates.append(match.group(0).decode("latin-1", errors="ignore"))

        cleaned = OrderedDict()
        for candidate in candidates:
            text = re.sub(r"\s+", " ", candidate).strip()
            if _looks_like_document_text(text):
                cleaned[text] = None

        content = "\n".join(cleaned.keys())
        if not content:
            raise RuntimeError(f"No readable text found in legacy {self.source_format} file.")

        return [
            Document(
                page_content=content,
                metadata={
                    "legacy_binary_extraction": "binary_strings",
                    "source_format": self.source_format,
                },
            )
        ]


def _looks_like_document_text(text: str) -> bool:
    if len(text) < 8:
        return False
    letters = sum(character.isalpha() for character in text)
    return letters >= 4 and letters / max(len(text), 1) >= 0.25


def _clean_text(text: str) -> str:
    return text.replace("\x00", "")
